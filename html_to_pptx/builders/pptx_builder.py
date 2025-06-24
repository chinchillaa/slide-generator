"""
PPTX Builder Module

Builds PowerPoint presentations from parsed slide data.
Supports both python-pptx library and native XML generation.
"""

import os
import sys
import zipfile
import tempfile
from typing import Dict, List, Any, Optional
from ..utils.xml_templates import XMLTemplates

# Try to import python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTXBuilder:
    """Build PPTX presentations from structured data"""
    
    def __init__(self, use_native: bool = False):
        """
        Initialize the builder
        
        Args:
            use_native: Force native XML generation even if python-pptx is available
        """
        self.use_native = use_native or not PPTX_AVAILABLE
        self.templates = XMLTemplates()
        
    def build(self, slides: List[Dict[str, Any]], output_path: str, metadata: Dict[str, str] = None):
        """
        Build PPTX file from slides data
        
        Args:
            slides: List of slide data dictionaries
            output_path: Path to save the PPTX file
            metadata: Document metadata (title, subtitle, etc.)
        """
        if self.use_native:
            self._build_native(slides, output_path, metadata)
        else:
            self._build_with_library(slides, output_path, metadata)
            
    def _build_with_library(self, slides: List[Dict[str, Any]], output_path: str, metadata: Dict[str, str] = None):
        """Build using python-pptx library"""
        prs = Presentation()
        
        # Set slide size to 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Process each slide
        for slide_data in slides:
            self._add_slide_with_library(prs, slide_data, metadata)
            
        # Save presentation
        prs.save(output_path)
        
    def _add_slide_with_library(self, prs, slide_data: Dict[str, Any], metadata: Dict[str, str] = None):
        """Add a single slide using python-pptx"""
        # Use blank layout
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Extract metadata from slide
        slide_meta = slide_data.get('metadata', {})
        
        # Add purple accent line
        if slide_meta.get('title'):
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), Inches(1.2),
                Inches(9), Inches(0.05)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(138, 43, 226)
            shape.line.fill.background()
        
        # Add title
        if slide_meta.get('title'):
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_meta['title']
            title_frame.paragraphs[0].font.size = Pt(36)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.name = 'Meiryo'
            
        # Add subtitle
        if slide_meta.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_meta['subtitle']
            subtitle_frame.paragraphs[0].font.size = Pt(28)
            subtitle_frame.paragraphs[0].font.bold = True
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(138, 43, 226)
            subtitle_frame.paragraphs[0].font.name = 'Meiryo'
            
        # Add period/additional info
        if slide_meta.get('period'):
            period_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
            period_frame = period_box.text_frame
            period_frame.text = f"対象期間: {slide_meta['period']}"
            period_frame.paragraphs[0].font.size = Pt(18)
            period_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
            period_frame.paragraphs[0].font.name = 'Meiryo'
            
        # Add logo
        logo_box = slide.shapes.add_textbox(Inches(8), Inches(0.3), Inches(1.5), Inches(0.5))
        logo_frame = logo_box.text_frame
        logo_frame.text = "= SB C&S"
        logo_frame.paragraphs[0].font.size = Pt(18)
        logo_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)
        logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add bottom gradient line
        gradient_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(5.3),
            Inches(10), Inches(0.05)
        )
        gradient_shape.fill.solid()
        gradient_shape.fill.fore_color.rgb = RGBColor(138, 43, 226)
        gradient_shape.line.fill.background()
        
    def _build_native(self, slides: List[Dict[str, Any]], output_path: str, metadata: Dict[str, str] = None):
        """Build using native XML generation"""
        with tempfile.TemporaryDirectory() as temp_dir:
            # Create directory structure
            self._create_directory_structure(temp_dir)
            
            # Create XML files
            self._create_content_types(temp_dir)
            self._create_relationships(temp_dir)
            self._create_document_properties(temp_dir, metadata)
            self._create_presentation(temp_dir, len(slides))
            self._create_theme(temp_dir)
            self._create_slide_master(temp_dir)
            self._create_slide_layout(temp_dir)
            
            # Create slides
            for idx, slide_data in enumerate(slides, 1):
                self._create_slide(temp_dir, idx, slide_data)
                
            # Package as PPTX
            self._create_pptx_package(temp_dir, output_path)
            
    def _create_directory_structure(self, base_path: str):
        """Create required directory structure"""
        dirs = [
            "_rels",
            "docProps",
            "ppt",
            "ppt/_rels",
            "ppt/slides",
            "ppt/slides/_rels",
            "ppt/slideLayouts",
            "ppt/slideLayouts/_rels",
            "ppt/slideMasters",
            "ppt/slideMasters/_rels",
            "ppt/theme"
        ]
        for dir_path in dirs:
            os.makedirs(os.path.join(base_path, dir_path), exist_ok=True)
            
    def _create_content_types(self, base_path: str):
        """Create [Content_Types].xml"""
        with open(os.path.join(base_path, "[Content_Types].xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_content_types())
            
    def _create_relationships(self, base_path: str):
        """Create relationship files"""
        # Root relationships
        with open(os.path.join(base_path, "_rels", ".rels"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_root_relationships())
            
        # Presentation relationships
        with open(os.path.join(base_path, "ppt", "_rels", "presentation.xml.rels"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_presentation_relationships(1))  # 1 slide for now
            
    def _create_document_properties(self, base_path: str, metadata: Dict[str, str] = None):
        """Create document properties"""
        metadata = metadata or {}
        
        # Core properties
        with open(os.path.join(base_path, "docProps", "core.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_core_properties(
                metadata.get('title', 'Presentation'),
                metadata.get('creator', 'HTML to PPTX Converter')
            ))
            
        # App properties
        with open(os.path.join(base_path, "docProps", "app.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_app_properties())
            
    def _create_presentation(self, base_path: str, slide_count: int):
        """Create presentation.xml"""
        with open(os.path.join(base_path, "ppt", "presentation.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_presentation(slide_count))
            
    def _create_theme(self, base_path: str):
        """Create theme file"""
        with open(os.path.join(base_path, "ppt", "theme", "theme1.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_theme())
            
    def _create_slide_master(self, base_path: str):
        """Create slide master"""
        with open(os.path.join(base_path, "ppt", "slideMasters", "slideMaster1.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_slide_master())
            
        # Master relationships
        with open(os.path.join(base_path, "ppt", "slideMasters", "_rels", "slideMaster1.xml.rels"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_slide_master_relationships())
            
    def _create_slide_layout(self, base_path: str):
        """Create slide layout"""
        with open(os.path.join(base_path, "ppt", "slideLayouts", "slideLayout1.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_slide_layout())
            
        # Layout relationships
        with open(os.path.join(base_path, "ppt", "slideLayouts", "_rels", "slideLayout1.xml.rels"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_slide_layout_relationships())
            
    def _create_slide(self, base_path: str, slide_num: int, slide_data: Dict[str, Any]):
        """Create individual slide"""
        slide_meta = slide_data.get('metadata', {})
        
        # Extract text content from elements if metadata is empty
        if not slide_meta.get('title') and slide_data.get('elements'):
            # Try to find title from h1 elements
            for element in slide_data['elements']:
                if hasattr(element, 'type') and element.type == 'h1' and hasattr(element, 'content'):
                    slide_meta['title'] = element.content
                    break
        
        # Create slide XML
        with open(os.path.join(base_path, "ppt", "slides", f"slide{slide_num}.xml"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_slide(
                slide_meta.get('title', 'Slide ' + str(slide_num)),
                slide_meta.get('subtitle', ''),
                slide_meta.get('period', '')
            ))
            
        # Create slide relationships
        with open(os.path.join(base_path, "ppt", "slides", "_rels", f"slide{slide_num}.xml.rels"), "w", encoding="utf-8") as f:
            f.write(self.templates.get_slide_relationships())
            
    def _create_pptx_package(self, temp_dir: str, output_path: str):
        """Package all files into PPTX"""
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arc_name = os.path.relpath(file_path, temp_dir)
                    zipf.write(file_path, arc_name)